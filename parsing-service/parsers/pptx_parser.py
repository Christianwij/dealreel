from pptx import Presentation
from loguru import logger
import io
import time
from typing import List, Dict, Any

from exceptions import (
    DocumentCorruptedError,
    TextExtractionError,
    TableExtractionError,
    MetadataExtractionError
)
from monitoring import record_extraction_metric, record_extraction_time

async def parse_pptx(file_content: bytes) -> dict:
    """Parse PPTX file content and return structured data."""
    try:
        # Create presentation from bytes
        pptx_stream = io.BytesIO(file_content)
        prs = Presentation(pptx_stream)
        
        # Extract metadata
        metadata_start = time.time()
        try:
            core_properties = prs.core_properties
            metadata = {
                "title": core_properties.title or "",
                "author": core_properties.author or "",
                "subject": core_properties.subject or "",
                "keywords": core_properties.keywords or "",
                "created": core_properties.created.isoformat() if core_properties.created else "",
                "modified": core_properties.modified.isoformat() if core_properties.modified else "",
            }
            record_extraction_time("metadata", "pptx", time.time() - metadata_start)
        except Exception as e:
            logger.error(f"Failed to extract PPTX metadata: {str(e)}")
            raise MetadataExtractionError(
                "Failed to extract PPTX metadata",
                document_type="pptx",
                details={"error": str(e)}
            )
        
        # Extract slides content
        slides = []
        text_content = []
        tables = []
        
        text_start = time.time()
        try:
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = {
                    "number": slide_num,
                    "shapes": [],
                    "notes": slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                }
                
                # Process shapes
                for shape in slide.shapes:
                    try:
                        if shape.has_text_frame:
                            text = shape.text.strip()
                            if text:
                                slide_content["shapes"].append({
                                    "type": "text",
                                    "content": text
                                })
                                text_content.append(text)
                                record_extraction_metric("text_shapes", "pptx")
                                
                        elif shape.has_table:
                            table_data = []
                            for row in shape.table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                if any(row_data):  # Skip empty rows
                                    table_data.append(row_data)
                            if table_data:
                                tables.append(table_data)
                                slide_content["shapes"].append({
                                    "type": "table",
                                    "content": table_data
                                })
                                record_extraction_metric("tables", "pptx")
                                
                        elif hasattr(shape, "image"):
                            try:
                                image_info = {
                                    "type": "image",
                                    "width": shape.width,
                                    "height": shape.height
                                }
                                slide_content["shapes"].append(image_info)
                                record_extraction_metric("images", "pptx")
                            except Exception as e:
                                logger.warning(f"Failed to extract image from slide {slide_num}: {str(e)}")
                                
                    except Exception as e:
                        logger.warning(f"Failed to process shape in slide {slide_num}: {str(e)}")
                        continue
                
                slides.append(slide_content)
            
            record_extraction_time("text", "pptx", time.time() - text_start)
            
        except Exception as e:
            logger.error(f"Failed to extract PPTX content: {str(e)}")
            raise TextExtractionError(
                "Failed to extract PPTX content",
                document_type="pptx",
                details={"error": str(e)}
            )
        
        return {
            "content_type": "pptx",
            "metadata": metadata,
            "text": "\n".join(text_content),
            "tables": tables,
            "slides": slides,
            "total_slides": len(prs.slides)
        }
        
    except (TextExtractionError, MetadataExtractionError) as e:
        raise e
    except Exception as e:
        logger.error(f"Failed to parse PPTX: {str(e)}")
        raise DocumentCorruptedError(
            "Failed to parse PPTX document",
            document_type="pptx",
            details={"error": str(e)}
        )
    finally:
        if 'pptx_stream' in locals():
            pptx_stream.close() 